"""
PPTX Builder - utilities for creating editable PPTX files
Based on OpenDCAI/DataFlow-Agent's implementation
"""
import os
import logging
from typing import List, Dict, Any, Optional, Tuple
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from PIL import Image

logger = logging.getLogger(__name__)


class PPTXBuilder:
    """Builder class for creating editable PPTX files from structured content"""
    
    # Standard slide dimensions (16:9 aspect ratio)
    DEFAULT_SLIDE_WIDTH_INCHES = 10
    DEFAULT_SLIDE_HEIGHT_INCHES = 5.625
    
    # Default DPI for pixel to inch conversion
    DEFAULT_DPI = 96
    
    # Font size ranges based on text level
    FONT_SIZE_RANGES = {
        1: (28, 72),  # Title/Heading 1 - allow larger titles
        2: (20, 48),  # Heading 2
        3: (16, 36),  # Heading 3
        'title': (28, 72),  # Title
        'default': (10, 32),  # Body text - wider range
        'header': (8, 16),  # Header
        'footer': (8, 16),  # Footer
    }
    
    def __init__(self, slide_width_inches: float = None, slide_height_inches: float = None):
        """
        Initialize PPTX builder
        
        Args:
            slide_width_inches: Slide width in inches (default: 10)
            slide_height_inches: Slide height in inches (default: 5.625)
        """
        self.slide_width_inches = slide_width_inches or self.DEFAULT_SLIDE_WIDTH_INCHES
        self.slide_height_inches = slide_height_inches or self.DEFAULT_SLIDE_HEIGHT_INCHES
        self.prs = None
        self.current_slide = None
        
    def create_presentation(self) -> Presentation:
        """Create a new presentation with configured dimensions"""
        self.prs = Presentation()
        self.prs.slide_width = Inches(self.slide_width_inches)
        self.prs.slide_height = Inches(self.slide_height_inches)
        return self.prs
    
    def setup_presentation_size(self, width_pixels: int, height_pixels: int, dpi: int = None):
        """
        Setup presentation size based on pixel dimensions
        
        Args:
            width_pixels: Width in pixels
            height_pixels: Height in pixels
            dpi: DPI for conversion (default: 96)
        """
        dpi = dpi or self.DEFAULT_DPI
        self.slide_width_inches = width_pixels / dpi
        self.slide_height_inches = height_pixels / dpi
        
        if self.prs:
            self.prs.slide_width = Inches(self.slide_width_inches)
            self.prs.slide_height = Inches(self.slide_height_inches)
    
    def add_blank_slide(self):
        """Add a blank slide to the presentation"""
        if not self.prs:
            self.create_presentation()
        
        # Use blank layout (layout 6 is typically blank)
        blank_layout = self.prs.slide_layouts[6]
        self.current_slide = self.prs.slides.add_slide(blank_layout)
        return self.current_slide
    
    def pixels_to_inches(self, pixels: float, dpi: int = None) -> float:
        """
        Convert pixels to inches
        
        Args:
            pixels: Pixel value
            dpi: DPI for conversion (default: 96)
            
        Returns:
            Value in inches
        """
        dpi = dpi or self.DEFAULT_DPI
        return pixels / dpi
    
    def calculate_font_size(self, bbox: List[int], text: str, text_level: Any = None, dpi: int = None) -> int:
        """
        Calculate appropriate font size based on bounding box and text
        Uses iterative approach to find the largest font that fits in the bbox
        
        Args:
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            text: Text content
            text_level: Text level (1=title, 2=heading, etc.) or type string
            dpi: DPI for pixel to inch conversion
            
        Returns:
            Font size in points
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Get bbox dimensions in pixels
        width_px = bbox[2] - bbox[0]
        height_px = bbox[3] - bbox[1]
        
        # Convert to inches for PPTX calculations
        width_in = width_px / dpi
        height_in = height_px / dpi
        
        # Determine font size range based on text level
        if isinstance(text_level, int):
            size_range = self.FONT_SIZE_RANGES.get(text_level, self.FONT_SIZE_RANGES['default'])
        elif isinstance(text_level, str):
            size_range = self.FONT_SIZE_RANGES.get(text_level, self.FONT_SIZE_RANGES['default'])
        else:
            size_range = self.FONT_SIZE_RANGES['default']
        
        min_size, max_size = size_range
        
        # Text metrics (approximate)
        text_length = len(text)
        
        # For very short text (likely titles or labels), use larger font
        if text_length <= 3:
            # Single characters or very short text
            estimated_size = int(height_in * 0.7 * 72)  # 72 points per inch, use 70% of height
            return max(min_size, min(max_size, estimated_size))
        
        # Estimate number of lines based on text length and box width
        # Assume average character width is about 0.6 * font_size in points
        # and we need to convert to the same unit system
        
        # Start with an estimate based on height
        # Typical line height is about 1.2x font size
        max_lines = max(1, int(height_in / 0.15))  # Assume minimum 0.15 inch per line
        
        # Binary search for optimal font size
        best_size = min_size
        
        for font_size in range(max_size, min_size - 1, -1):
            # Estimate how many characters fit per line
            # Average character width for proportional fonts is about 0.5-0.6 of font size
            char_width_pts = font_size * 0.55  # points
            char_width_in = char_width_pts / 72  # inches
            
            # Account for padding (roughly 5% on each side)
            usable_width = width_in * 0.9
            chars_per_line = int(usable_width / char_width_in)
            
            if chars_per_line < 1:
                continue
            
            # Calculate required lines
            required_lines = max(1, (text_length + chars_per_line - 1) // chars_per_line)
            
            # Line height is typically 1.2x font size
            line_height_pts = font_size * 1.2
            line_height_in = line_height_pts / 72
            
            total_height_needed = required_lines * line_height_in
            
            # Add some padding (10%)
            total_height_needed *= 1.1
            
            # If text fits, this is our font size
            if total_height_needed <= height_in:
                best_size = font_size
                break
        
        return best_size
    
    def add_text_element(
        self,
        slide,
        text: str,
        bbox: List[int],
        text_level: Any = None,
        dpi: int = None,
        align: str = 'left'
    ):
        """
        Add text element to slide
        
        Args:
            slide: Target slide
            text: Text content
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            text_level: Text level (1=title, 2=heading, etc.) or type string
            dpi: DPI for conversion (default: 96)
            align: Text alignment ('left', 'center', 'right')
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Convert bbox to inches
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        
        # Add text box
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = text
        text_frame.word_wrap = True
        
        # Set font size (pass original bbox in pixels and dpi)
        font_size = self.calculate_font_size(bbox, text, text_level, dpi)
        paragraph = text_frame.paragraphs[0]
        paragraph.font.size = Pt(font_size)
        
        # Remove default margins for better fit
        text_frame.margin_left = Inches(0.05)
        text_frame.margin_right = Inches(0.05)
        text_frame.margin_top = Inches(0.05)
        text_frame.margin_bottom = Inches(0.05)
        
        # Set alignment
        if align == 'center':
            paragraph.alignment = PP_ALIGN.CENTER
        elif align == 'right':
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = PP_ALIGN.LEFT
        
        # Make title text bold
        if text_level == 1 or text_level == 'title':
            paragraph.font.bold = True
        
        # Calculate bbox dimensions for logging
        bbox_width = bbox[2] - bbox[0]
        bbox_height = bbox[3] - bbox[1]
        size_range = self.FONT_SIZE_RANGES.get(text_level, self.FONT_SIZE_RANGES['default'])
        logger.debug(f"Text: '{text[:35]}' | box: {bbox_width}x{bbox_height}px | level: {text_level} (range: {size_range[0]}-{size_range[1]}pt) | font: {font_size}pt")
    
    def add_image_element(
        self,
        slide,
        image_path: str,
        bbox: List[int],
        dpi: int = None
    ):
        """
        Add image element to slide
        
        Args:
            slide: Target slide
            image_path: Path to image file
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            dpi: DPI for conversion (default: 96)
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Check if image exists
        if not os.path.exists(image_path):
            logger.warning(f"Image not found: {image_path}, adding placeholder")
            self.add_image_placeholder(slide, bbox, dpi)
            return
        
        # Convert bbox to inches
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        
        try:
            # Add image
            slide.shapes.add_picture(image_path, left, top, width, height)
            logger.debug(f"Added image: {image_path} at bbox {bbox}")
        except Exception as e:
            logger.error(f"Failed to add image {image_path}: {str(e)}")
            self.add_image_placeholder(slide, bbox, dpi)
    
    def add_image_placeholder(
        self,
        slide,
        bbox: List[int],
        dpi: int = None
    ):
        """
        Add a placeholder for missing images
        
        Args:
            slide: Target slide
            bbox: Bounding box [x0, y0, x1, y1] in pixels
            dpi: DPI for conversion (default: 96)
        """
        dpi = dpi or self.DEFAULT_DPI
        
        # Convert bbox to inches
        left = Inches(self.pixels_to_inches(bbox[0], dpi))
        top = Inches(self.pixels_to_inches(bbox[1], dpi))
        width = Inches(self.pixels_to_inches(bbox[2] - bbox[0], dpi))
        height = Inches(self.pixels_to_inches(bbox[3] - bbox[1], dpi))
        
        # Add a text box as placeholder
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.text = "[Image]"
        paragraph = text_frame.paragraphs[0]
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(12)
        paragraph.font.italic = True
    
    def save(self, output_path: str):
        """
        Save presentation to file
        
        Args:
            output_path: Output file path
        """
        if not self.prs:
            raise ValueError("No presentation to save")
        
        # Ensure directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        self.prs.save(output_path)
        logger.info(f"Saved presentation to: {output_path}")
    
    def get_presentation(self) -> Presentation:
        """Get the current presentation object"""
        return self.prs

